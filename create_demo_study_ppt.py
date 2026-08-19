"""
Generate a study/demo PowerPoint for the Healthcare Analytics Dashboard.
Run: python create_demo_study_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from copy import deepcopy
import os

# Colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
SKY = RGBColor(0x60, 0xA5, 0xFA)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
CARD = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)

OUT = os.path.join(os.path.dirname(__file__), "Healthcare_Dashboard_Demo_Study_PPT.pptx")


def set_run(run, size=18, bold=False, color=DARK_TEXT, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK_TEXT, spacing=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if spacing:
            p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)
    return box


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_bar(slide, 0, Inches(6.9), Inches(13.333), Inches(0.6), BLUE)

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "MCA Project  |  Demo Day Study Deck", size=16, color=SKY, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.4),
                "Healthcare Data Analysis &\nVisualization Dashboard",
                size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8),
                "Easy study slides — project overview, architecture, pages, ML & code",
                size=18, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
                "Streamlit  •  Python  •  SQLite  •  Plotly  •  scikit-learn\n10,000 patient records  |  6 analytics pages  |  3 ML features",
                size=16, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "Use these slides to revise + speak confidently to your lecturer",
                size=13, color=WHITE)


def section_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.5),
                f"SECTION  {number}", size=14, bold=True, color=SKY)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.0),
                title, size=34, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.6),
                    subtitle, size=18, color=MUTED)
    return slide


def content_header(slide, title, subtitle=""):
    add_bg(slide, LIGHT_BG)
    add_bar(slide, 0, 0, Inches(13.333), Inches(1.05), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.28), Inches(12), Inches(0.5),
                title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
                    subtitle, size=14, color=GRAY)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Study Agenda", "What you will learn from this deck")
    items = [
        ("01", "Project Pitch & Problem"),
        ("02", "Tech Stack & How to Run"),
        ("03", "Architecture & Data Flow"),
        ("04", "Dataset & Database"),
        ("05", "All 6 Dashboard Pages"),
        ("06", "Machine Learning Explained"),
        ("07", "Code Structure (Easy)"),
        ("08", "Demo Script & Viva Tips"),
    ]
    for i, (num, text) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.9 + row * 1.2)
        add_card(slide, left, top, Inches(5.8), Inches(1.0))
        add_textbox(slide, left + Inches(0.25), top + Inches(0.28), Inches(0.8), Inches(0.5),
                    num, size=22, bold=True, color=BLUE)
        add_textbox(slide, left + Inches(1.1), top + Inches(0.32), Inches(4.4), Inches(0.5),
                    text, size=18, bold=True, color=DARK_TEXT)


def slide_pitch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "30-Second Elevator Pitch", "Memorize this — open your demo with it")
    add_card(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8))
    quote = (
        "Sir/Madam, this is a Healthcare Analytics Dashboard.\n\n"
        "We took a Kaggle dataset of 10,000 patient records, cleaned it, "
        "stored it in SQLite, and built a multi-page Streamlit app.\n\n"
        "It shows demographics, disease trends, hospital operations, "
        "financial analysis, and Machine Learning — K-Means clustering, "
        "Random Forest prediction, and patient risk scoring.\n\n"
        "Everything runs locally with three commands: install, load data, run Streamlit."
    )
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.2),
                quote, size=17, color=DARK_TEXT)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Problem → Solution", "Why this project exists")

    add_card(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.8))
    add_bar(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(0.55), RED)
    add_textbox(slide, Inches(0.85), Inches(1.82), Inches(5.3), Inches(0.4),
                "WITHOUT Dashboard", size=16, bold=True, color=WHITE)
    add_bullets(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.6), [
        "Scattered CSV / Excel files",
        "Manual chart making",
        "Hard to find trends quickly",
        "No prediction or risk view",
        "Management decisions are slow",
    ], size=16)

    add_card(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.8))
    add_bar(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(0.55), GREEN)
    add_textbox(slide, Inches(7.15), Inches(1.82), Inches(5.3), Inches(0.4),
                "WITH Our Dashboard", size=16, bold=True, color=WHITE)
    add_bullets(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.6), [
        "One clean SQLite database",
        "Auto interactive Plotly charts",
        "Instant KPIs + filters",
        "ML clustering + prediction + risk",
        "Faster hospital insights",
    ], size=16)


def slide_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Technology Stack", "Say these names confidently in viva")

    rows = [
        ("UI Framework", "Streamlit", "Python web dashboard — no HTML/JS needed"),
        ("Language", "Python", "Main coding language of the project"),
        ("Database", "SQLite + SQLAlchemy", "Single file DB — easy to demo anywhere"),
        ("Data Processing", "Pandas + NumPy", "Cleaning, grouping, calculations"),
        ("Charts", "Plotly", "Interactive graphs (hover, zoom)"),
        ("Machine Learning", "scikit-learn", "K-Means, Random Forest, PCA"),
        ("Dataset", "Kaggle healthcare-dataset", "10,000 real public-domain records"),
    ]
    y = 1.35
    for layer, tech, meaning in rows:
        add_card(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.72))
        add_textbox(slide, Inches(0.7), Inches(y + 0.18), Inches(2.6), Inches(0.4),
                    layer, size=13, bold=True, color=GRAY)
        add_textbox(slide, Inches(3.4), Inches(y + 0.18), Inches(3.4), Inches(0.4),
                    tech, size=15, bold=True, color=BLUE)
        add_textbox(slide, Inches(6.9), Inches(y + 0.18), Inches(5.6), Inches(0.4),
                    meaning, size=13, color=DARK_TEXT)
        y += 0.8


def slide_run(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "How to Run (3 Steps)", "Practice this before demo day")

    steps = [
        ("1", "Install", "pip install -r requirements.txt", "Installs Streamlit, pandas, plotly, sklearn…"),
        ("2", "Load Data", "python database/loader.py", "Creates SQLite DB + loads 10,000 records"),
        ("3", "Launch App", "streamlit run app.py", "Opens dashboard at localhost:8501"),
    ]
    for i, (num, title, cmd, desc) in enumerate(steps):
        left = Inches(0.55 + i * 4.2)
        add_card(slide, left, Inches(1.8), Inches(3.95), Inches(4.4))
        # circle number
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.45), Inches(2.15), Inches(1.0), Inches(1.0))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        add_textbox(slide, left + Inches(1.45), Inches(2.35), Inches(1.0), Inches(0.6),
                    num, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(3.4), Inches(3.55), Inches(0.4),
                    title, size=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(3.95), Inches(3.55), Inches(0.7),
                    cmd, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.25), Inches(4.8), Inches(3.45), Inches(1.0),
                    desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Project Architecture (Easy Map)", "Think like a restaurant kitchen")

    layers = [
        ("CSV / Kaggle", "Raw ingredients"),
        ("loader.py", "Kitchen — clean & store"),
        ("analytics/", "Chefs — SQL queries"),
        ("visualizations/", "Plating — charts/KPIs"),
        ("pages/ + app.py", "Menu — what user sees"),
    ]
    for i, (name, meaning) in enumerate(layers):
        top = Inches(1.45 + i * 1.05)
        add_card(slide, Inches(2.5), top, Inches(8.3), Inches(0.9))
        add_bar(slide, Inches(2.5), top, Inches(0.18), Inches(0.9), BLUE)
        add_textbox(slide, Inches(3.0), top + Inches(0.12), Inches(4.5), Inches(0.35),
                    name, size=18, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(3.0), top + Inches(0.48), Inches(7.2), Inches(0.3),
                    meaning, size=13, color=GRAY)
        if i < len(layers) - 1:
            add_textbox(slide, Inches(6.2), top + Inches(0.78), Inches(1), Inches(0.3),
                        "↓", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def slide_folders(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Folder Structure", "What each folder does")

    left_items = [
        "app.py → Home page",
        "pages/ → 6 dashboard screens",
        "analytics/ → SQL business logic",
        "visualizations/ → charts + KPI cards",
    ]
    right_items = [
        "database/ → connection, schema, loader",
        "data/ → CSV + healthcare.db",
        "assets/style.css → dark theme",
        "config/settings.py → colors/config",
    ]
    add_card(slide, Inches(0.6), Inches(1.6), Inches(6.0), Inches(4.9))
    add_textbox(slide, Inches(0.9), Inches(1.85), Inches(5.4), Inches(0.4),
                "UI & Logic", size=18, bold=True, color=BLUE)
    add_bullets(slide, Inches(0.9), Inches(2.5), Inches(5.4), Inches(3.6), left_items, size=16)

    add_card(slide, Inches(6.9), Inches(1.6), Inches(6.0), Inches(4.9))
    add_textbox(slide, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.4),
                "Data & Config", size=18, bold=True, color=BLUE)
    add_bullets(slide, Inches(7.2), Inches(2.5), Inches(5.4), Inches(3.6), right_items, size=16)


def slide_dataflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Data Flow", "Follow the path of one patient record")

    steps = [
        "CSV / Kaggle\nor Synthetic",
        "loader.py\nclean + enrich",
        "SQLite\nhealthcare.db",
        "analytics/\nSQL queries",
        "pages/\ncharts + KPIs",
    ]
    for i, text in enumerate(steps):
        left = Inches(0.4 + i * 2.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.6), Inches(2.25), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        add_textbox(slide, left + Inches(0.1), Inches(3.1), Inches(2.05), Inches(1.0),
                    text, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, left + Inches(2.15), Inches(3.2), Inches(0.5), Inches(0.5),
                        "→", size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.7))
    add_bullets(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.3), [
        "loader.py also creates age_group and length_of_stay (extra useful columns)",
        "Every page calls analytics functions → query_df(SQL) → Pandas DataFrame → Plotly chart",
        "Layered design: UI does not hold heavy SQL — analytics modules do",
    ], size=14)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Dataset Snapshot", "What is inside the 10,000 records?")

    cards = [
        ("10,000", "Patient records"),
        ("15+", "Original CSV columns"),
        ("1", "Main SQLite table"),
        ("CC0", "Public domain license"),
    ]
    for i, (big, small) in enumerate(cards):
        left = Inches(0.55 + i * 3.2)
        add_card(slide, left, Inches(1.55), Inches(3.0), Inches(1.7))
        add_textbox(slide, left + Inches(0.1), Inches(1.75), Inches(2.8), Inches(0.7),
                    big, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(2.55), Inches(2.8), Inches(0.4),
                    small, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    groups = [
        ("Patient", "Name, Age, Gender, Blood Type, Age Group"),
        ("Clinical", "Medical Condition, Medication, Test Results"),
        ("Operations", "Hospital, Doctor, Room, Admission Type, LOS"),
        ("Financial", "Billing Amount, Insurance Provider"),
    ]
    for i, (title, cols) in enumerate(groups):
        top = Inches(3.55 + i * 0.85)
        add_card(slide, Inches(0.55), top, Inches(12.2), Inches(0.75))
        add_textbox(slide, Inches(0.8), top + Inches(0.2), Inches(2.2), Inches(0.4),
                    title, size=15, bold=True, color=BLUE)
        add_textbox(slide, Inches(3.2), top + Inches(0.2), Inches(9.2), Inches(0.4),
                    cols, size=14, color=DARK_TEXT)


def slide_pages_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Dashboard Pages (Map)", "6 pages + Home — click order for demo")

    pages = [
        ("01", "Overview", "KPIs, monthly trends, recent records"),
        ("02", "Demographics", "Age, gender, blood type, insurance"),
        ("03", "Disease", "Conditions, meds, test results"),
        ("04", "Operations", "Hospital/doctor, LOS, admissions"),
        ("05", "Financial", "Revenue, payer mix, high-value"),
        ("06", "ML Insights", "Clustering, RF, risk scoring"),
    ]
    for i, (num, name, desc) in enumerate(pages):
        col = i % 3
        row = i // 3
        left = Inches(0.55 + col * 4.2)
        top = Inches(1.6 + row * 2.6)
        add_card(slide, left, top, Inches(3.95), Inches(2.3))
        add_textbox(slide, left + Inches(0.25), top + Inches(0.35), Inches(3.4), Inches(0.45),
                    num, size=24, bold=True, color=BLUE)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.9), Inches(3.4), Inches(0.4),
                    name, size=18, bold=True, color=DARK_TEXT)
        add_textbox(slide, left + Inches(0.25), top + Inches(1.4), Inches(3.4), Inches(0.7),
                    desc, size=13, color=GRAY)


def add_page_detail(prs, num, title, say, shows, code):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, f"Page {num} — {title}", "What to show + what to say")

    add_card(slide, Inches(0.55), Inches(1.55), Inches(12.2), Inches(1.5))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.35),
                "SAY THIS", size=12, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.7),
                say, size=15, color=DARK_TEXT)

    add_card(slide, Inches(0.55), Inches(3.3), Inches(6.0), Inches(3.4))
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.4),
                "What it shows", size=16, bold=True, color=DARK_TEXT)
    add_bullets(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.4), shows, size=14)

    add_card(slide, Inches(6.8), Inches(3.3), Inches(5.95), Inches(3.4))
    add_textbox(slide, Inches(7.05), Inches(3.5), Inches(5.5), Inches(0.4),
                "Code file", size=16, bold=True, color=DARK_TEXT)
    add_textbox(slide, Inches(7.05), Inches(4.2), Inches(5.5), Inches(2.0),
                code, size=14, color=GRAY)


def slide_ml_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "ML Insights — Overview", "Most impressive page — spend 2–3 minutes here")

    cards = [
        ("K-Means", "Unsupervised", "Group similar patients\ninto cost/stay segments", GREEN),
        ("Random Forest", "Supervised", "Predict test result:\nNormal / Abnormal / Inconclusive", BLUE),
        ("Risk Score", "Rule-based", "Score patients 0–100\nLow / Moderate / High", AMBER),
    ]
    for i, (name, typ, desc, color) in enumerate(cards):
        left = Inches(0.55 + i * 4.2)
        add_card(slide, left, Inches(1.7), Inches(3.95), Inches(4.6))
        add_bar(slide, left, Inches(1.7), Inches(3.95), Inches(0.15), color)
        add_textbox(slide, left + Inches(0.25), Inches(2.2), Inches(3.45), Inches(0.5),
                    name, size=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.25), Inches(2.8), Inches(3.45), Inches(0.4),
                    typ, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.3), Inches(3.6), Inches(3.35), Inches(1.5),
                    desc, size=15, color=GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.25), Inches(5.4), Inches(3.45), Inches(0.5),
                    "File: analytics/ml_insights.py", size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_ml_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "ML — How Each Model Works", "Simple explanations for viva")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.7))
    add_textbox(slide, Inches(0.75), Inches(1.65), Inches(11.8), Inches(0.35),
                "K-Means + PCA", size=15, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.75), Inches(2.1), Inches(11.8), Inches(0.9),
                "Features: age, billing, LOS, gender, admission type, condition → StandardScaler → "
                "K-Means (2–8 clusters) → PCA projects to 2D for scatter plot. "
                "Why scale? Billing numbers are huge vs age — without scaling, billing dominates.",
                size=13, color=DARK_TEXT)

    add_card(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.7))
    add_textbox(slide, Inches(0.75), Inches(3.55), Inches(11.8), Inches(0.35),
                "Random Forest Classifier", size=15, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.9),
                "Label = test_results. Train/test = 80/20. 100 trees vote. "
                "Shows accuracy, confusion matrix, feature importance. "
                "Why RF? Strong ensemble model, handles mixed features, easy to explain.",
                size=13, color=DARK_TEXT)

    add_card(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6))
    add_textbox(slide, Inches(0.75), Inches(5.45), Inches(11.8), Inches(0.35),
                "Risk Score Formula", size=15, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.75), Inches(5.9), Inches(11.8), Inches(0.7),
                "0.30×Age + 0.30×Billing + 0.25×LOS + 0.10×AbnormalTest + 0.05×Emergency → score 0–100 "
                "→ Low (0–33) / Moderate (34–66) / High (67–100)",
                size=13, color=DARK_TEXT)


def slide_code(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Code Pattern (All Pages Same)", "Understand once — explain every page")

    steps = [
        ("1", "Page setup", "set_page_config + load CSS theme"),
        ("2", "Import logic", "from analytics.* import get_..."),
        ("3", "Fetch data", "SQL via query_df() → DataFrame"),
        ("4", "Show KPIs", "metric_row([...]) big number cards"),
        ("5", "Draw charts", "Plotly helpers + st.plotly_chart"),
        ("6", "Tables", "st.dataframe for detail rows"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        left = Inches(0.55 + col * 4.2)
        top = Inches(1.6 + row * 2.55)
        add_card(slide, left, top, Inches(3.95), Inches(2.25))
        add_textbox(slide, left + Inches(0.25), top + Inches(0.35), Inches(3.4), Inches(0.45),
                    num, size=28, bold=True, color=BLUE)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.95), Inches(3.4), Inches(0.4),
                    title, size=17, bold=True, color=DARK_TEXT)
        add_textbox(slide, left + Inches(0.25), top + Inches(1.45), Inches(3.4), Inches(0.55),
                    desc, size=13, color=GRAY)


def slide_key_files(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Key Files to Open During Demo", "Show 4 files — enough to impress")

    files = [
        ("database/loader.py", "How data enters the system\n(Kaggle → CSV → clean → SQLite)"),
        ("database/connection.py", "query_df() — SQL becomes\nPandas DataFrame"),
        ("analytics/overview.py", "Example SQL analytics\n(KPIs, monthly trends)"),
        ("analytics/ml_insights.py", "K-Means + Random Forest\n+ Risk scoring code"),
    ]
    for i, (path, desc) in enumerate(files):
        col = i % 2
        row = i // 2
        left = Inches(0.55 + col * 6.4)
        top = Inches(1.6 + row * 2.6)
        add_card(slide, left, top, Inches(6.1), Inches(2.3))
        add_textbox(slide, left + Inches(0.35), top + Inches(0.45), Inches(5.4), Inches(0.5),
                    path, size=16, bold=True, color=BLUE)
        add_textbox(slide, left + Inches(0.35), top + Inches(1.15), Inches(5.4), Inches(0.8),
                    desc, size=14, color=DARK_TEXT)


def slide_demo_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "8–10 Minute Demo Flow", "Follow this order on project day")

    flow = [
        ("0:45", "Intro pitch"),
        ("0:30", "Home + DB connected"),
        ("1:00", "Overview KPIs"),
        ("1:00", "Demographics"),
        ("1:00", "Disease page"),
        ("1:00", "Operations"),
        ("1:00", "Financial"),
        ("2:30", "ML Insights ★"),
        ("1:00", "Show 2–3 code files"),
        ("0:20", "Closing line"),
    ]
    for i, (time, step) in enumerate(flow):
        col = i % 5
        row = i // 5
        left = Inches(0.4 + col * 2.55)
        top = Inches(1.7 + row * 2.5)
        add_card(slide, left, top, Inches(2.4), Inches(2.1))
        add_textbox(slide, left + Inches(0.1), top + Inches(0.45), Inches(2.2), Inches(0.45),
                    time, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(1.1), Inches(2.2), Inches(0.7),
                    step, size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def slide_viva(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "Likely Viva Questions", "Short ready answers")

    qas = [
        ("What is your project?", "Healthcare analytics dashboard on 10k records with SQL + charts + ML."),
        ("Why Streamlit?", "Build interactive dashboard in pure Python — fast for data projects."),
        ("Why SQLite?", "Zero setup, one file, portable, enough for 10k rows."),
        ("Clustering vs Classification?", "Clustering = no labels (groups). Classification = predict labeled class."),
        ("How is LOS calculated?", "discharge_date − date_of_admission (days) in loader.py."),
        ("Architecture pattern?", "Layered: database → analytics → visualizations → pages."),
    ]
    y = 1.4
    for q, a in qas:
        add_textbox(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(0.3),
                    f"Q: {q}", size=13, bold=True, color=BLUE)
        add_textbox(slide, Inches(0.6), Inches(y + 0.32), Inches(12.1), Inches(0.35),
                    f"A: {a}", size=13, color=DARK_TEXT)
        y += 0.9


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_bar(slide, 0, Inches(6.9), Inches(13.333), Inches(0.6), BLUE)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5),
                "CLOSING LINE (memorize)", size=14, bold=True, color=SKY)
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(2.2),
                "Our project follows a clean layered design:\ndata is loaded once into SQLite,\nanalytics modules run SQL,\nvisualization helpers draw charts,\nand Streamlit pages present insights —\nincluding clustering, prediction, and risk scoring.",
                size=20, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.8),
                "Study tip: revise DEMO_STUDY_NOTES.md + practice this PPT once with the live app open.",
                size=14, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(7.05), Inches(11.5), Inches(0.35),
                "Thank you  |  Ready for questions", size=14, bold=True, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    slide_agenda(prs)

    section_slide(prs, "01", "Project Pitch & Problem", "Start strong — explain purpose clearly")
    slide_pitch(prs)
    slide_problem(prs)

    section_slide(prs, "02", "Tech Stack & Setup", "Tools + how to launch the app")
    slide_stack(prs)
    slide_run(prs)

    section_slide(prs, "03", "Architecture & Data", "How folders and data connect")
    slide_architecture(prs)
    slide_folders(prs)
    slide_dataflow(prs)
    slide_dataset(prs)

    section_slide(prs, "04", "Dashboard Pages", "What each screen means for your demo")
    slide_pages_overview(prs)

    add_page_detail(
        prs, "01", "Overview",
        "This is the executive summary for management — one screen to understand the whole hospital dataset.",
        ["Total records, patients, hospitals, doctors", "Revenue, avg billing, avg LOS",
         "Abnormal test rate", "Monthly admissions & revenue trends", "Recent records table"],
        "pages/01_Overview.py\n+\nanalytics/overview.py\n\nKey functions:\nget_kpis()\nget_monthly_trend()",
    )
    add_page_detail(
        prs, "02", "Patient Demographics",
        "Who are our patients? Age, gender, blood type, insurance mix — useful for planning.",
        ["Age group distribution", "Male / Female breakdown", "Blood type chart",
         "Insurance provider mix", "Age × gender heatmap", "Age vs billing scatter"],
        "pages/02_Patient_Demographics.py\n+\nanalytics/patient_analytics.py",
    )
    add_page_detail(
        prs, "03", "Disease Analysis",
        "Which diseases are most common, which medicines are used, and how test results vary.",
        ["Top medical conditions", "Medication summary", "Test results by condition",
         "Condition trends over time", "Treemap / severity proxy", "Top-N slider filter"],
        "pages/03_Disease_Analysis.py\n+\nanalytics/disease_analytics.py",
    )
    add_page_detail(
        prs, "04", "Hospital Operations",
        "Operations view — hospital/doctor load, length of stay, emergency vs elective.",
        ["Hospital performance table", "Doctor workload (top 50)", "Admission type mix",
         "LOS distribution", "Weekly admission pattern", "Room utilization"],
        "pages/04_Hospital_Operations.py\n+\nanalytics/hospital_analytics.py",
    )
    add_page_detail(
        prs, "05", "Financial Analysis",
        "Financial view — where revenue comes from, insurance mix, and expensive cases.",
        ["Total / avg / min / max billing", "Revenue by condition & hospital",
         "Payer mix (insurance)", "Monthly revenue trend", "High-value cases filter"],
        "pages/05_Financial_Analysis.py\n+\nanalytics/financial_analytics.py\n\nSidebar: high-value threshold",
    )

    section_slide(prs, "05", "Machine Learning", "The page that impresses most")
    slide_ml_intro(prs)
    slide_ml_detail(prs)

    section_slide(prs, "06", "Code & Demo Day", "Explain code simply + follow the script")
    slide_code(prs)
    slide_key_files(prs)
    slide_demo_flow(prs)
    slide_viva(prs)
    slide_closing(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
